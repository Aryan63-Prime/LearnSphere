"""
LearnSphere — Export Utilities
Generate downloadable PPT presentations from AI-generated content.
Uses python-pptx for PPTX generation.
"""

import os
import re
import uuid
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ── Paths ──────────────────────────────────────────────────────
EXPORT_DIR = Path(__file__).parent / 'data' / 'exports'
EXPORT_DIR.mkdir(parents=True, exist_ok=True)

# ── Color Palette (matches Obsidian Ember) ─────────────────────
BG_DARK     = RGBColor(0x1B, 0x13, 0x2B)
BG_CARD     = RGBColor(0x27, 0x1D, 0x3A)
GOLD        = RGBColor(0xF5, 0xC5, 0x42)
CORAL       = RGBColor(0xFF, 0x6B, 0x6B)
TEAL        = RGBColor(0x4E, 0xCE, 0xC4)
LAVENDER    = RGBColor(0xC0, 0x92, 0xFD)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF0, 0xEB, 0xE3)
MUTED       = RGBColor(0xAA, 0x9E, 0xBC)
DIM         = RGBColor(0x7B, 0x72, 0x8E)

TITLE_FONT = 'Georgia'
BODY_FONT  = 'Segoe UI'


# ── Helpers ────────────────────────────────────────────────────

def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text,
              sz=18, color=WHITE, bold=False, font=BODY_FONT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tf


def _add_para(tf, text, sz=14, color=MUTED, bold=False, font=BODY_FONT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.space_before = Pt(4)
    p.space_after = Pt(4)
    return p


def _accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


# ── Markdown → Sections Parser ─────────────────────────────────

def _parse_markdown_sections(markdown_text: str) -> list[dict]:
    """
    Parse markdown into structured sections for slides.
    Each section has a 'title' and 'bullets' list.
    """
    sections = []
    current = {'title': 'Overview', 'bullets': []}

    for line in markdown_text.split('\n'):
        line = line.strip()
        if not line:
            continue

        # Heading → new section
        if line.startswith('#'):
            if current['bullets']:
                sections.append(current)
            title = re.sub(r'^#+\s*', '', line).strip()
            current = {'title': title, 'bullets': []}
        # Bullet point
        elif line.startswith(('- ', '* ', '• ')):
            bullet = re.sub(r'^[-*•]\s*', '', line)
            # Strip bold markdown
            bullet = re.sub(r'\*\*(.*?)\*\*', r'\1', bullet)
            current['bullets'].append(bullet)
        # Numbered list
        elif re.match(r'^\d+[.)]\s', line):
            bullet = re.sub(r'^\d+[.)]\s*', '', line)
            bullet = re.sub(r'\*\*(.*?)\*\*', r'\1', bullet)
            current['bullets'].append(bullet)
        # Plain text → add as bullet
        else:
            clean = re.sub(r'\*\*(.*?)\*\*', r'\1', line)
            current['bullets'].append(clean)

    if current['bullets']:
        sections.append(current)

    return sections if sections else [{'title': 'Content', 'bullets': [markdown_text[:500]]}]


# ── PPT Builder ────────────────────────────────────────────────

def generate_ppt(topic: str, content: str, modality: str = 'text') -> dict:
    """
    Generate a PPTX presentation from AI-generated content.

    Args:
        topic:     The subject matter
        content:   Markdown/text content to export
        modality:  Which modality generated this ('text', 'code', etc.)

    Returns:
        dict with success, filename, filepath
    """
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        SLIDE_W = prs.slide_width
        SLIDE_H = prs.slide_height

        accent_colors = [GOLD, CORAL, TEAL, LAVENDER]

        # ── Slide 1: Title ──────────────────────────────────
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(s, BG_DARK)
        _accent_bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GOLD)
        _accent_bar(s, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), LAVENDER)

        _add_text(s, Inches(1.5), Inches(1.5), Inches(10), Inches(1.0),
                  'LearnSphere', sz=52, color=OFF_WHITE, bold=True, font=TITLE_FONT)

        _add_text(s, Inches(1.5), Inches(2.7), Inches(10), Inches(0.8),
                  topic, sz=32, color=GOLD, bold=False, font=TITLE_FONT)

        _accent_bar(s, Inches(1.5), Inches(3.8), Inches(4), Inches(0.03), GOLD)

        modality_labels = {
            'text': '📖 Text Explanation',
            'code': '⚡ Code Generation',
            'audio': '🎵 Audio Script',
            'visual': '🖼️ Visual Diagram',
        }
        _add_text(s, Inches(1.5), Inches(4.2), Inches(8), Inches(0.5),
                  modality_labels.get(modality, '📖 Content'), sz=18, color=MUTED)

        _add_text(s, Inches(1.5), Inches(5.5), Inches(8), Inches(0.5),
                  'Generated by LearnSphere AI • Hackaholics Squad', sz=14, color=DIM)

        # ── Content Slides ──────────────────────────────────
        sections = _parse_markdown_sections(content)

        for idx, section in enumerate(sections):
            s = prs.slides.add_slide(prs.slide_layouts[6])
            _set_bg(s, BG_DARK)

            accent = accent_colors[idx % len(accent_colors)]

            # Top accent + header
            _accent_bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.05), accent)
            _add_text(s, Inches(1.0), Inches(0.4), Inches(11), Inches(0.8),
                      section['title'], sz=34, color=OFF_WHITE, bold=True, font=TITLE_FONT)
            _accent_bar(s, Inches(1.0), Inches(1.2), Inches(2.5), Inches(0.04), accent)

            # Content bullets
            if section['bullets']:
                tf = _add_text(s, Inches(1.0), Inches(1.6), Inches(11.3), Inches(5.3),
                               section['bullets'][0], sz=15, color=MUTED, font=BODY_FONT)
                for bullet in section['bullets'][1:]:
                    _add_para(tf, bullet, sz=15, color=MUTED)

        # ── Save ────────────────────────────────────────────
        safe_topic = re.sub(r'[^\w\s-]', '', topic)[:40].strip().replace(' ', '_')
        filename = f"LearnSphere_{safe_topic}_{uuid.uuid4().hex[:6]}.pptx"
        filepath = EXPORT_DIR / filename
        prs.save(str(filepath))

        return {
            'success': True,
            'filename': filename,
            'filepath': str(filepath),
            'num_slides': len(prs.slides),
        }

    except Exception as e:
        return {'success': False, 'error': str(e)}
