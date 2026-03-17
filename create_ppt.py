"""
LearnSphere PPT Generator — Warm Sunset Edition
A unique, premium presentation with warm gradients, serif + sans-serif typography,
diagonal accent shapes, and a distinctive editorial magazine-style layout.
Team: Hackaholics Squad
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Color Palette: "Sunset Noir" ──────────────────────────────
BG_RICH      = RGBColor(0x1B, 0x13, 0x2B)    # Deep indigo-black
BG_CARD      = RGBColor(0x27, 0x1D, 0x3A)    # Dark plum card
WARM_GOLD    = RGBColor(0xF5, 0xC5, 0x42)    # Warm gold
SUNSET_CORAL = RGBColor(0xFF, 0x6B, 0x6B)    # Coral red
MINT_TEAL    = RGBColor(0x4E, 0xCE, 0xC4)    # Teal mint
LAVENDER     = RGBColor(0xC0, 0x92, 0xFD)    # Soft lavender
SOFT_PEACH   = RGBColor(0xFF, 0xA0, 0x7A)    # Light salmon/peach
ELECTRIC_LIME = RGBColor(0xA8, 0xE6, 0x4E)   # Lime green
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE    = RGBColor(0xF0, 0xEB, 0xE3)    # Cream white
MUTED_TEXT   = RGBColor(0xAA, 0x9E, 0xBC)    # Muted lavender-gray
DIM_TEXT     = RGBColor(0x7B, 0x72, 0x8E)

# Fonts
TITLE_FONT  = 'Georgia'              # Elegant serif for titles
BODY_FONT   = 'Segoe UI'             # Clean sans-serif for body
CODE_FONT   = 'Cascadia Code'        # Monospace for code/structure

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ── Helpers ────────────────────────────────────────────────────
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, border=None, bw=Pt(0), radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = bw
    else:
        shape.line.fill.background()
    return shape


def flat_rect(slide, l, t, w, h, fill, border=None, bw=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = bw
    else:
        shape.line.fill.background()
    return shape


def txt(slide, l, t, w, h, text, sz=18, clr=WHITE, bold=False,
        align=PP_ALIGN.LEFT, font=BODY_FONT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = clr
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def add_para(tf, text, sz=14, clr=MUTED_TEXT, bold=False, align=PP_ALIGN.LEFT,
             font=BODY_FONT, sb=Pt(3), sa=Pt(3)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = clr
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.space_before = sb
    p.space_after = sa
    return p


def accent_line(slide, l, t, w, h, clr):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = clr
    bar.line.fill.background()
    return bar


def diamond(slide, l, t, size, clr):
    d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, l, t, size, size)
    d.fill.solid()
    d.fill.fore_color.rgb = clr
    d.line.fill.background()
    return d


def circle(slide, l, t, size, clr, text="", tsz=20):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    c.fill.solid()
    c.fill.fore_color.rgb = clr
    c.line.fill.background()
    if text:
        tf = c.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(tsz)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    return c


def slide_header(slide, title, accent_clr, subtitle=None):
    """Reusable slide header with diagonal accent."""
    # Top accent stripe
    accent_line(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), accent_clr)
    # Decorative diamond
    diamond(slide, Inches(0.6), Inches(0.5), Inches(0.35), accent_clr)
    # Title
    txt(slide, Inches(1.2), Inches(0.35), Inches(10), Inches(0.8),
        title, sz=38, clr=OFF_WHITE, bold=True, font=TITLE_FONT)
    # Underline
    accent_line(slide, Inches(1.2), Inches(1.2), Inches(2.2), Inches(0.04), accent_clr)
    if subtitle:
        txt(slide, Inches(1.2), Inches(1.35), Inches(10), Inches(0.5),
            subtitle, sz=15, clr=DIM_TEXT, font=BODY_FONT)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_RICH)

# Decorative top gold bar
accent_line(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), WARM_GOLD)

# Large vertical accent strip on the left
flat_rect(s, Inches(0), Inches(0), Inches(0.6), SLIDE_H, RGBColor(0x22, 0x18, 0x33))
accent_line(s, Inches(0.55), Inches(1.0), Inches(0.05), Inches(5.5), WARM_GOLD)

# Decorative scattered diamonds
diamond(s, Inches(10.5), Inches(1.0), Inches(0.6), RGBColor(0x35, 0x28, 0x4A))
diamond(s, Inches(11.4), Inches(1.8), Inches(0.4), RGBColor(0x30, 0x24, 0x45))
diamond(s, Inches(10.0), Inches(5.5), Inches(0.5), RGBColor(0x2D, 0x22, 0x42))

# Small gold diamonds as decoration
diamond(s, Inches(11.0), Inches(3.0), Inches(0.25), WARM_GOLD)
diamond(s, Inches(11.8), Inches(4.5), Inches(0.2), SUNSET_CORAL)
diamond(s, Inches(10.3), Inches(4.0), Inches(0.15), MINT_TEAL)

# Project emoji badge
circle(s, Inches(1.2), Inches(1.2), Inches(1.3), RGBColor(0x35, 0x28, 0x4A), "🎓", tsz=42)

# Title text
txt(s, Inches(3.0), Inches(1.2), Inches(8), Inches(1.0),
    "LearnSphere", sz=60, clr=OFF_WHITE, bold=True, font=TITLE_FONT)

# Subtitle
txt(s, Inches(3.0), Inches(2.3), Inches(8), Inches(0.7),
    "AI-Powered Learning Platform", sz=26, clr=WARM_GOLD, bold=False, font=TITLE_FONT)

# Small description
txt(s, Inches(3.0), Inches(3.2), Inches(7), Inches(0.5),
    "Multi-modal educational experience powered by Google Gemini 2.0 Flash",
    sz=16, clr=MUTED_TEXT)

# Gold separator
accent_line(s, Inches(3.0), Inches(4.2), Inches(4.0), Inches(0.03), WARM_GOLD)

# Team name
txt(s, Inches(3.0), Inches(4.6), Inches(6), Inches(0.6),
    "HACKAHOLICS SQUAD", sz=22, clr=SUNSET_CORAL, bold=True, font=BODY_FONT)

# Members
txt(s, Inches(3.0), Inches(5.3), Inches(8), Inches(0.5),
    "Aryan Patel   ·   Chetan Singh   ·   Tejas Singh Nirvan",
    sz=16, clr=MUTED_TEXT)

# Bottom accent
accent_line(s, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), LAVENDER)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_RICH)
slide_header(s, "Problem Statement", SUNSET_CORAL, "Why we built LearnSphere")

problems = [
    ("01", "Rigid Learning Formats",
     "Traditional platforms deliver content in a single format — text only. "
     "Visual, auditory, and kinesthetic learners are underserved.",
     SUNSET_CORAL),
    ("02", "Fragmented Tool Ecosystem",
     "Students juggle between multiple disconnected apps for explanations, "
     "code examples, audio notes, and visual diagrams.",
     WARM_GOLD),
    ("03", "No On-Demand Personalization",
     "Existing tools lack AI-driven depth control. Learners can't choose "
     "between a quick overview and a deep technical dive.",
     LAVENDER),
]

for i, (num, title, desc, clr) in enumerate(problems):
    top = Inches(1.8 + i * 1.8)

    # Number badge
    circle(s, Inches(0.8), top + Inches(0.15), Inches(0.7), clr, num, tsz=18)

    # Card
    rect(s, Inches(1.8), top, Inches(10.8), Inches(1.5), BG_CARD, clr, Pt(1))

    txt(s, Inches(2.1), top + Inches(0.15), Inches(10), Inches(0.45),
        title, sz=21, clr=clr, bold=True, font=TITLE_FONT)

    txt(s, Inches(2.1), top + Inches(0.65), Inches(10.2), Inches(0.7),
        desc, sz=14, clr=MUTED_TEXT)

# Solution callout at bottom
rect(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.9), RGBColor(0x1A, 0x2E, 0x27), MINT_TEAL, Pt(1.5))
txt(s, Inches(1.2), Inches(6.3), Inches(11), Inches(0.7),
    "💡  Our Answer:  A single unified platform that adapts to every learning style — "
    "text, code, audio & visuals — all powered by AI.",
    sz=15, clr=MINT_TEAL, bold=True)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — Core Features
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_RICH)
slide_header(s, "Core Features", WARM_GOLD, "Four powerful learning modalities in one platform")

features = [
    ("📖", "Text Explanation",
     "AI-generated explanations with\n3 depth levels: Brief,\nComprehensive & In-Depth.\nMarkdown-formatted output\nwith examples & applications.",
     WARM_GOLD),
    ("⚡", "Code Generation",
     "Production-ready Python code\nwith inline documentation,\ndependency lists & usage\nexamples. Supports Basic,\nIntermediate & Advanced.",
     MINT_TEAL),
    ("🎵", "Audio Learning",
     "Natural speech scripts\noptimized for TTS playback.\nChoose Conversational,\nFormal or Storytelling\nnarration styles.",
     SUNSET_CORAL),
    ("🖼️", "Visual Diagrams",
     "Auto-generated Mermaid.js\nflowcharts & architecture\ndiagrams. Real-time rendering\nwith styled, color-coded\nnodes & connections.",
     LAVENDER),
]

for i, (icon, title, desc, clr) in enumerate(features):
    left = Inches(0.5 + i * 3.15)
    top  = Inches(1.8)

    # Card with thick left border effect
    rect(s, left, top, Inches(2.9), Inches(5.0), BG_CARD, clr, Pt(1))
    accent_line(s, left, top, Inches(0.06), Inches(5.0), clr)

    # Icon
    txt(s, left + Inches(0.2), top + Inches(0.3), Inches(2.5), Inches(0.7),
        icon, sz=40, align=PP_ALIGN.LEFT)

    # Title
    txt(s, left + Inches(0.2), top + Inches(1.1), Inches(2.5), Inches(0.5),
        title, sz=20, clr=clr, bold=True, font=TITLE_FONT)

    # Thin separator
    accent_line(s, left + Inches(0.2), top + Inches(1.7), Inches(1.5), Inches(0.02), clr)

    # Desc
    txt(s, left + Inches(0.2), top + Inches(1.9), Inches(2.5), Inches(2.8),
        desc, sz=13, clr=MUTED_TEXT)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — Technology Stack
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_RICH)
slide_header(s, "Technology Stack", MINT_TEAL, "Built with modern, robust technologies")

techs = [
    ("🐍", "Backend",     "Python 3.9+\nFlask Framework\nRESTful JSON APIs",     WARM_GOLD),
    ("🎨", "Frontend",    "HTML5 · CSS3 · JS\nJinja2 Templates\nGlassmorphism UI", LAVENDER),
    ("🤖", "AI Engine",   "Google Gemini 2.0 Flash\nPerplexity Sonar (Fallback)\nCustom Prompts",     MINT_TEAL),
    ("🔊", "Audio",       "Web Speech API\nText-to-Speech\n3 Narration Styles",   SUNSET_CORAL),
    ("📊", "Diagrams",    "Mermaid.js\nDynamic Flowcharts\nReal-time Render",     ELECTRIC_LIME),
]

# Staggered cards layout
for i, (icon, label, items, clr) in enumerate(techs):
    left = Inches(0.35 + i * 2.55)
    top  = Inches(1.9 + (0.3 if i % 2 == 1 else 0))

    rect(s, left, top, Inches(2.35), Inches(4.5), BG_CARD, clr, Pt(1))

    # Top colored strip
    accent_line(s, left, top, Inches(2.35), Inches(0.06), clr)

    # Icon + Label
    txt(s, left + Inches(0.15), top + Inches(0.25), Inches(2.0), Inches(0.6),
        f"{icon}", sz=34, align=PP_ALIGN.CENTER)

    txt(s, left + Inches(0.1), top + Inches(1.0), Inches(2.15), Inches(0.5),
        label, sz=19, clr=clr, bold=True, align=PP_ALIGN.CENTER, font=TITLE_FONT)

    accent_line(s, left + Inches(0.4), top + Inches(1.6), Inches(1.55), Inches(0.02), clr)

    txt(s, left + Inches(0.1), top + Inches(1.8), Inches(2.15), Inches(2.2),
        items, sz=13, clr=MUTED_TEXT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — System Architecture
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_RICH)
slide_header(s, "System Architecture", LAVENDER, "Clean layered design from user to AI")

layers = [
    ("🌐  Browser UI",     "HTML5 Templates  ·  Glassmorphism CSS  ·  JavaScript Client", WARM_GOLD),
    ("⚡  Flask Server",   "app.py  —  Route Handling  ·  Request Parsing  ·  JSON Response",   SUNSET_CORAL),
    ("🧠  AI Utility Layer", "genai_utils.py  —  Prompt Builder  ·  API Orchestration  ·  Parsing", MINT_TEAL),
    ("🤖  AI Models",      "Google Gemini 2.0 Flash  ·  Perplexity Sonar (Fallback)",         LAVENDER),
]

for i, (label, desc, clr) in enumerate(layers):
    top = Inches(1.8 + i * 1.35)

    # Wide layer card
    rect(s, Inches(1.0), top, Inches(11.3), Inches(1.05), BG_CARD, clr, Pt(1.5))

    # Left color stripe
    accent_line(s, Inches(1.0), top, Inches(0.08), Inches(1.05), clr)

    txt(s, Inches(1.4), top + Inches(0.1), Inches(5), Inches(0.4),
        label, sz=19, clr=clr, bold=True, font=TITLE_FONT)

    txt(s, Inches(1.4), top + Inches(0.55), Inches(10.5), Inches(0.4),
        desc, sz=13, clr=MUTED_TEXT)

    # Arrow connectors
    if i < 3:
        arrow_top = top + Inches(1.08)
        # Vertical connector line
        accent_line(s, Inches(6.6), arrow_top, Inches(0.04), Inches(0.27), DIM_TEXT)
        # Small diamond at connection point
        diamond(s, Inches(6.5), arrow_top + Inches(0.08), Inches(0.2), clr)

# Right side callout box
rect(s, Inches(8.5), Inches(6.3), Inches(4.3), Inches(0.8), BG_CARD, WARM_GOLD, Pt(1))
txt(s, Inches(8.7), Inches(6.4), Inches(3.9), Inches(0.6),
    "⚙️  API Config: Temp 0.7 · Max 4000 tokens · 60s timeout",
    sz=12, clr=WARM_GOLD, bold=True)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — Project Structure
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_RICH)
slide_header(s, "Project Structure", ELECTRIC_LIME, "Clean, modular codebase organization")

file_tree = """LearnSphere/
├── app.py                 ← Flask Application
├── config.py              ← Config Manager
├── genai_utils.py         ← AI Integration
├── requirements.txt       ← Dependencies
├── .env                   ← API Keys (secured)
├── templates/
│   ├── base.html          ← Base Layout
│   ├── index.html         ← Home Page
│   ├── text_explanation.html
│   ├── code_generation.html
│   ├── audio_learning.html
│   └── image_visualization.html
└── static/
    ├── css/style.css      ← UI Styling
    └── js/main.js         ← Client Logic"""

# Code card
rect(s, Inches(0.5), Inches(1.7), Inches(6.8), Inches(5.3), BG_CARD, ELECTRIC_LIME, Pt(1))
accent_line(s, Inches(0.5), Inches(1.7), Inches(6.8), Inches(0.05), ELECTRIC_LIME)

txt(s, Inches(0.8), Inches(1.95), Inches(6.3), Inches(4.8),
    file_tree, sz=12, clr=ELECTRIC_LIME, font=CODE_FONT)

# Key files - right side
files = [
    ("app.py", "Flask app with 5 page routes and\n4 JSON API endpoints for each\nlearning modality.", WARM_GOLD),
    ("genai_utils.py", "211 lines of AI logic: custom\nprompt engineering for each mode,\nGemini & Perplexity API calls.", MINT_TEAL),
    ("config.py", "Centralized configuration from\n.env: API keys, debug flags,\nand Flask secret key.", LAVENDER),
    ("templates/", "6 Jinja2 HTML templates with\nglassmorphism design, responsive\nlayout & interactive elements.", SUNSET_CORAL),
]

for i, (fname, desc, clr) in enumerate(files):
    top = Inches(1.7 + i * 1.3)
    rect(s, Inches(7.8), top, Inches(5.0), Inches(1.15), BG_CARD, clr, Pt(1))
    accent_line(s, Inches(7.8), top, Inches(0.06), Inches(1.15), clr)

    txt(s, Inches(8.1), top + Inches(0.1), Inches(4.5), Inches(0.3),
        fname, sz=15, clr=clr, bold=True, font=CODE_FONT)

    txt(s, Inches(8.1), top + Inches(0.42), Inches(4.5), Inches(0.65),
        desc, sz=11, clr=MUTED_TEXT)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — How It Works
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_RICH)
slide_header(s, "How It Works", WARM_GOLD, "Simple 5-step user journey")

steps = [
    ("①", "Choose\nMode",    "Pick Text, Code,\nAudio or Visual",    WARM_GOLD),
    ("②", "Enter\nTopic",    "Type any subject\nyou want to learn",  SUNSET_CORAL),
    ("③", "Set\nPreferences", "Adjust depth,\ncomplexity or style", LAVENDER),
    ("④", "AI\nGenerates",   "Gemini 2.0 Flash\nprocesses your query", MINT_TEAL),
    ("⑤", "Learn &\nExport",  "Read, listen, copy\nor download results",  ELECTRIC_LIME),
]

y_center = Inches(3.8)

for i, (num, title, desc, clr) in enumerate(steps):
    left = Inches(0.4 + i * 2.6)

    # Circle with number
    circle(s, left + Inches(0.55), Inches(1.8), Inches(1.2), clr, num, tsz=30)

    # Title
    txt(s, left + Inches(0.1), Inches(3.3), Inches(2.1), Inches(0.8),
        title, sz=18, clr=clr, bold=True, align=PP_ALIGN.CENTER, font=TITLE_FONT)

    # Card
    rect(s, left, Inches(4.3), Inches(2.3), Inches(1.6), BG_CARD, clr, Pt(1))
    accent_line(s, left, Inches(4.3), Inches(2.3), Inches(0.04), clr)

    txt(s, left + Inches(0.15), Inches(4.5), Inches(2.0), Inches(1.2),
        desc, sz=12, clr=MUTED_TEXT, align=PP_ALIGN.CENTER)

    # Connector arrows
    if i < 4:
        accent_line(s, left + Inches(2.3), Inches(2.35), Inches(0.3), Inches(0.03), DIM_TEXT)
        diamond(s, left + Inches(2.35), Inches(2.28), Inches(0.15), clr)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — AI & Prompt Engineering
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_RICH)
slide_header(s, "AI & Prompt Engineering", SUNSET_CORAL, "Smart prompt design for each learning modality")

# Left panel - Model specs
rect(s, Inches(0.5), Inches(1.8), Inches(5.8), Inches(5.2), BG_CARD, SUNSET_CORAL, Pt(1))
accent_line(s, Inches(0.5), Inches(1.8), Inches(5.8), Inches(0.05), SUNSET_CORAL)

txt(s, Inches(0.8), Inches(2.0), Inches(5.2), Inches(0.5),
    "🤖  Gemini 2.0 Flash Configuration", sz=20, clr=SUNSET_CORAL, bold=True, font=TITLE_FONT)

specs = [
    ("API Endpoint",   "generativelanguage.googleapis.com"),
    ("Temperature",    "0.7 — balanced creativity & accuracy"),
    ("Max Tokens",     "4,000 tokens per request"),
    ("Timeout",        "60 seconds"),
    ("Fallback Model", "Perplexity Sonar (REST API)"),
]

for i, (key, val) in enumerate(specs):
    top = Inches(2.7 + i * 0.55)
    txt(s, Inches(0.9), top, Inches(2.2), Inches(0.4),
        key, sz=13, clr=WARM_GOLD, bold=True)
    txt(s, Inches(3.1), top, Inches(3.0), Inches(0.4),
        val, sz=13, clr=MUTED_TEXT)

# Prompt modes
accent_line(s, Inches(0.9), Inches(5.5), Inches(4.8), Inches(0.02), DIM_TEXT)
txt(s, Inches(0.9), Inches(5.65), Inches(5.0), Inches(0.4),
    "Prompt Modes", sz=16, clr=WARM_GOLD, bold=True, font=TITLE_FONT)

modes = [
    ("Text →", "Depth-adapted educational prompts", WARM_GOLD),
    ("Code →", "Complexity-leveled code generation", MINT_TEAL),
    ("Audio →", "Style-specific narration scripts", SUNSET_CORAL),
    ("Visual →", "Mermaid.js diagram instructions", LAVENDER),
]

for i, (mode, desc, clr) in enumerate(modes):
    row = i // 2
    col = i % 2
    left = Inches(0.9 + col * 2.7)
    top  = Inches(6.1 + row * 0.4)
    txt(s, left, top, Inches(1.0), Inches(0.35), mode, sz=11, clr=clr, bold=True)
    txt(s, left + Inches(0.9), top, Inches(2.0), Inches(0.35), desc, sz=11, clr=MUTED_TEXT)

# Right panel - API Flow
rect(s, Inches(6.8), Inches(1.8), Inches(6.0), Inches(5.2), BG_CARD, LAVENDER, Pt(1))
accent_line(s, Inches(6.8), Inches(1.8), Inches(6.0), Inches(0.05), LAVENDER)

txt(s, Inches(7.1), Inches(2.0), Inches(5.4), Inches(0.5),
    "⚡  Request → Response Pipeline", sz=20, clr=LAVENDER, bold=True, font=TITLE_FONT)

flow = [
    ("01", "User submits topic + preferences via browser",  WARM_GOLD),
    ("02", "Flask route receives POST with JSON body",      MUTED_TEXT),
    ("03", "genai_utils.py selects & builds custom prompt",  MINT_TEAL),
    ("04", "HTTP POST to Gemini 2.0 Flash API",             SUNSET_CORAL),
    ("05", "Response JSON parsed → content extracted",       MUTED_TEXT),
    ("06", "Metadata added (model, time, success flag)",     LAVENDER),
    ("07", "JSON result returned → rendered in browser",     ELECTRIC_LIME),
]

for i, (num, step, clr) in enumerate(flow):
    top = Inches(2.7 + i * 0.58)
    circle(s, Inches(7.2), top, Inches(0.35), clr, num, tsz=10)
    txt(s, Inches(7.7), top + Inches(0.02), Inches(4.8), Inches(0.35),
        step, sz=13, clr=MUTED_TEXT)


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — UI / UX Design
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_RICH)
slide_header(s, "UI / UX Design", LAVENDER, "Premium, modern interface design principles")

ui_items = [
    ("🌊", "Glassmorphism",
     "Frosted glass card effects with backdrop-filter blur "
     "and semi-transparent layers. Creates depth and elegance.",
     WARM_GOLD),
    ("🌙", "Dark-First Theme",
     "Deep indigo-black backgrounds with carefully tuned contrast "
     "ratios. Reduces eye strain during long study sessions.",
     LAVENDER),
    ("📱", "Responsive Design",
     "Flexbox grid adapts to all screen sizes. Desktop, tablet, "
     "and mobile optimized. No breakpoint issues.",
     MINT_TEAL),
    ("✨", "Micro-Interactions",
     "Smooth CSS transitions on hover. Loading spinners during AI "
     "generation. Fade-in content reveals for engagement.",
     SUNSET_CORAL),
    ("🎯", "Intuitive Navigation",
     "Persistent top navbar with clear feature links. Each learning "
     "mode is one click away from the home page.",
     ELECTRIC_LIME),
    ("📋", "Smart Output Rendering",
     "Markdown rendering for text, syntax highlighting for code, "
     "TTS controls for audio, live Mermaid diagrams for visuals.",
     SOFT_PEACH),
]

for i, (icon, title, desc, clr) in enumerate(ui_items):
    col = i % 3
    row = i // 3

    left = Inches(0.4 + col * 4.2)
    top  = Inches(1.7 + row * 2.7)

    rect(s, left, top, Inches(3.9), Inches(2.4), BG_CARD, clr, Pt(1))
    accent_line(s, left, top, Inches(3.9), Inches(0.04), clr)

    txt(s, left + Inches(0.2), top + Inches(0.2), Inches(1.0), Inches(0.5),
        icon, sz=28, align=PP_ALIGN.LEFT)

    txt(s, left + Inches(0.9), top + Inches(0.2), Inches(2.8), Inches(0.4),
        title, sz=18, clr=clr, bold=True, font=TITLE_FONT)

    txt(s, left + Inches(0.2), top + Inches(0.8), Inches(3.5), Inches(1.4),
        desc, sz=12, clr=MUTED_TEXT)


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — Future Scope
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_RICH)
slide_header(s, "Future Scope", SOFT_PEACH, "Planned enhancements and growth roadmap")

roadmap = [
    ("🔐", "User Authentication",    "Login/signup & personalized\nlearning profile tracking",     WARM_GOLD),
    ("📊", "Learning Analytics",     "Progress dashboard, streaks,\nand learning habit insights",   LAVENDER),
    ("🤝", "Collaborative Rooms",    "Real-time group study with\nshared AI-generated content",    MINT_TEAL),
    ("📄", "Export to PDF",          "Download explanations, code\nand diagrams as PDF files",      SUNSET_CORAL),
    ("🧪", "Quiz Generation",        "AI-created assessments for\nself-evaluation & testing",      ELECTRIC_LIME),
    ("🌍", "Multi-Language",         "Content generation in 10+\nlanguages for global access",     SOFT_PEACH),
]

for i, (icon, title, desc, clr) in enumerate(roadmap):
    col = i % 3
    row = i // 3

    left = Inches(0.4 + col * 4.2)
    top  = Inches(1.7 + row * 2.8)

    rect(s, left, top, Inches(3.9), Inches(2.4), BG_CARD, clr, Pt(1))

    # Left thick accent
    accent_line(s, left, top, Inches(0.07), Inches(2.4), clr)

    txt(s, left + Inches(0.3), top + Inches(0.2), Inches(3.3), Inches(0.5),
        f"{icon}  {title}", sz=19, clr=clr, bold=True, font=TITLE_FONT)

    accent_line(s, left + Inches(0.3), top + Inches(0.75), Inches(2.0), Inches(0.02), clr)

    txt(s, left + Inches(0.3), top + Inches(0.9), Inches(3.3), Inches(1.2),
        desc, sz=13, clr=MUTED_TEXT)


# ══════════════════════════════════════════════════════════════
# SLIDE 11 — Thank You
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_RICH)

# Top bar
accent_line(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), WARM_GOLD)

# Left accent strip
flat_rect(s, Inches(0), Inches(0), Inches(0.6), SLIDE_H, RGBColor(0x22, 0x18, 0x33))
accent_line(s, Inches(0.55), Inches(1.0), Inches(0.05), Inches(5.5), WARM_GOLD)

# Decorative diamonds
diamond(s, Inches(10.2), Inches(1.5), Inches(0.4), RGBColor(0x35, 0x28, 0x4A))
diamond(s, Inches(11.0), Inches(2.3), Inches(0.25), WARM_GOLD)
diamond(s, Inches(10.8), Inches(5.0), Inches(0.3), SUNSET_CORAL)
diamond(s, Inches(11.5), Inches(3.5), Inches(0.2), MINT_TEAL)

# Thank you text
txt(s, Inches(1.5), Inches(1.3), Inches(9), Inches(1.0),
    "Thank You", sz=56, clr=OFF_WHITE, bold=True, font=TITLE_FONT)

txt(s, Inches(1.5), Inches(2.4), Inches(9), Inches(0.6),
    "We appreciate your time and attention.", sz=22, clr=WARM_GOLD, font=TITLE_FONT)

# Gold separator
accent_line(s, Inches(1.5), Inches(3.3), Inches(4.0), Inches(0.03), WARM_GOLD)

# Team banner
txt(s, Inches(1.5), Inches(3.7), Inches(6), Inches(0.6),
    "HACKAHOLICS SQUAD", sz=24, clr=SUNSET_CORAL, bold=True, font=BODY_FONT)

# Member cards
members = [
    ("Aryan Patel",         WARM_GOLD),
    ("Chetan Singh",        MINT_TEAL),
    ("Tejas Singh Nirvan",  LAVENDER),
]

for i, (name, clr) in enumerate(members):
    left = Inches(1.5 + i * 3.2)
    top  = Inches(4.5)

    rect(s, left, top, Inches(2.8), Inches(1.6), BG_CARD, clr, Pt(1.5))
    accent_line(s, left, top, Inches(2.8), Inches(0.04), clr)

    circle(s, left + Inches(0.9), top + Inches(0.2), Inches(0.9), clr, "👤", tsz=24)

    txt(s, left + Inches(0.1), top + Inches(1.15), Inches(2.6), Inches(0.35),
        name, sz=14, clr=OFF_WHITE, bold=True, align=PP_ALIGN.CENTER, font=BODY_FONT)

# Footer
accent_line(s, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), LAVENDER)

txt(s, Inches(1.5), Inches(6.5), Inches(10), Inches(0.4),
    "🎓  LearnSphere  ·  AI-Powered Learning Platform  ·  Google Gemini 2.0 Flash",
    sz=13, clr=DIM_TEXT)


# ── Save ───────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "LearnSphere_Presentation_v2.pptx")
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
